import os
from pptx import Presentation
import json

def extract_text_from_pptx(file_path):
    """Extract text content from PowerPoint slides."""
    prs = Presentation(file_path)
    slides_content = []
    
    for i, slide in enumerate(prs.slides):
        slide_content = {
            "slide_number": i + 1,
            "slide_title": "",
            "slide_content": []
        }
        
        # Extract slide title if available
        if slide.shapes.title and slide.shapes.title.text:
            slide_content["slide_title"] = slide.shapes.title.text
        
        # Extract text from all shapes in the slide
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                slide_content["slide_content"].append(shape.text.strip())
        
        slides_content.append(slide_content)
    
    return slides_content

def analyze_pitch_decks(directory_path, output_dir):
    """Analyze all PPTX files in the given directory."""
    results = {}
    
    for filename in os.listdir(directory_path):
        if filename.endswith(".pptx"):
            file_path = os.path.join(directory_path, filename)
            print(f"Analyzing: {filename}")
            
            try:
                slides_content = extract_text_from_pptx(file_path)
                results[filename] = slides_content
                
                # Save individual deck analysis
                output_file = os.path.join(output_dir, f"{os.path.splitext(filename)[0]}_analysis.json")
                with open(output_file, 'w') as f:
                    json.dump(slides_content, f, indent=2)
                
                print(f"Saved analysis to: {output_file}")
                
            except Exception as e:
                print(f"Error processing {filename}: {str(e)}")
    
    # Save combined analysis
    combined_output = os.path.join(output_dir, "combined_analysis.json")
    with open(combined_output, 'w') as f:
        json.dump(results, f, indent=2)
    
    print(f"Saved combined analysis to: {combined_output}")
    
    return results

if __name__ == "__main__":
    upload_dir = "/home/ubuntu/upload"
    output_dir = "/home/ubuntu/jade_kite_analysis"
    analyze_pitch_decks(upload_dir, output_dir)
